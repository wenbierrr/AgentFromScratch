import os
from google import genai
from google.genai import types
from dotenv import load_dotenv
from eval import judge_summary_quality, evaluate

load_dotenv()
client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])

INPUT_FOLDER = "input"
OUTPUT_FOLDER = "output"
MAX_RETRIES = 3
SUPPORTED_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".py", ".yaml", ".yml", ".pdf", ".docx", ".pptx"}

os.makedirs(INPUT_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


# ─────────────────────────────────────────────
# TOOLS
# ─────────────────────────────────────────────

def read_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                text = "\n\n".join(page.extract_text() or "" for page in pdf.pages)
            return text if text.strip() else "ERROR: No readable text — may be scanned/image-only."

        elif ext == ".docx":
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext in {".pptx", ".ppt"}:
            from pptx import Presentation
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides) if slides else "ERROR: No readable text found in presentation."

        else:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()

    except Exception as e:
        return f"ERROR: Could not read file — {e}"


def markdown_to_blocks(text: str) -> list:
    blocks = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            blocks.append({"object": "block", "type": "heading_1",
                "heading_1": {"rich_text": [{"type": "text", "text": {"content": line[2:]}}]}})
        elif line.startswith("## "):
            blocks.append({"object": "block", "type": "heading_2",
                "heading_2": {"rich_text": [{"type": "text", "text": {"content": line[3:]}}]}})
        elif line.startswith("- ") or line.startswith("* "):
            blocks.append({"object": "block", "type": "bulleted_list_item",
                "bulleted_list_item": {"rich_text": [{"type": "text", "text": {"content": line[2:]}}]}})
        else:
            blocks.append({"object": "block", "type": "paragraph",
                "paragraph": {"rich_text": [{"type": "text", "text": {"content": line}}]}})
    return blocks


def write_to_notion(title: str, summary: str) -> str:
    from notion_client import Client
    notion = Client(auth=os.environ["NOTION_TOKEN"])
    database_id = os.environ["NOTION_DATABASE_ID"]

    try:
        blocks = markdown_to_blocks(summary)
        notion.pages.create(
            parent={"database_id": database_id},
            properties={
                "Name": {"title": [{"text": {"content": title}}]},
                "Summary": {"rich_text": [{"text": {"content": summary[:2000]}}]},
            },
            children=blocks[:100],
        )
        return f"Successfully saved '{title}' to Notion."
    except Exception as e:
        return f"ERROR: Could not write to Notion — {e}"


# ─────────────────────────────────────────────
# TOOL DECLARATIONS
# ─────────────────────────────────────────────

read_file_declaration = {
    "name": "read_file",
    "description": "Reads a file from local disk and returns its full contents. "
                   "Use this to see a file's contents before summarizing it.",
    "parameters": {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to the file, e.g. 'input/notes.txt'"},
        },
        "required": ["path"],
    },
}

write_to_notion_declaration = {
    "name": "write_to_notion",
    "description": "Saves a completed summary to the Notion database as a new page. "
                   "Always call this AFTER you have generated the summary — never before.",
    "parameters": {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Title for the page, e.g. the source filename without extension"},
            "summary": {"type": "string", "description": "The full markdown summary to save"},
        },
        "required": ["title", "summary"],
    },
}

available_tools = {"read_file": read_file, "write_to_notion": write_to_notion}

# ─────────────────────────────────────────────
# SYSTEM PROMPT
# ─────────────────────────────────────────────

SYSTEM_PROMPT = """You are a meticulous technical summarizer. You act as an intelligent filter: \
you condense dense material into its absolute core essentials so that even a non-technical reader grasps the main \
takeaway at a glance — without losing crucial context. You operate as an \
agent with tools that let you read source documents. The summaries you produce are saved as \
study notes the reader will return to later, so they must be accurate and self-contained.

YOUR TASK
When given a file or document:
1. Use the read_file tool to read the FULL source before summarizing. Never summarize a file you haven't read.
2. Produce the summary as a markdown document with these sections:
   - **Bottom line**: one sentence — the single most important takeaway.
   - **Key points**: 3-5 bullets, one distinct idea each, in order of importance.
   - **Common misconceptions**: Identify highly prevalent misunderstandings that could potentially arise (omit this section if there are none).
3. After generating the summary, call write_to_notion to save it. Do not skip this step.

RULES
- Be concise. Cut filler, keep substance. Use plain language, but keep the author's key terms where technical jargon matters.
- Preserve nuance. Never flatten a trade-off into a one-sided claim.
- Ground everything in the source. Do NOT add facts, examples, or opinions that aren't in the original.
- Mark uncertainty. If something is context-dependent or the author is speculating, signal it ("the author argues...").
- If you can't read the file or it's empty, say so plainly and stop. Do not guess at the contents."""

config = types.GenerateContentConfig(
    system_instruction=SYSTEM_PROMPT,
    tools=[types.Tool(function_declarations=[read_file_declaration, write_to_notion_declaration])],
    automatic_function_calling=types.AutomaticFunctionCallingConfig(disable=True),
)

# ─────────────────────────────────────────────
# MAIN: SCAN + PROCESS
# ─────────────────────────────────────────────

files = sorted([
    f for f in os.listdir(INPUT_FOLDER)
    if os.path.isfile(os.path.join(INPUT_FOLDER, f))
    and os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS
])

if not files:
    print(f"No supported files found in '{INPUT_FOLDER}/'. Drop a file in and try again.")
else:
    print(f"Found {len(files)} file(s): {files}\n")

for filename in files:
    filepath = os.path.join(INPUT_FOLDER, filename)
    print(f"\n{'='*50}\nProcessing: {filename}\n{'='*50}")

    source_text = ""
    summary_text = ""
    retry_count = 0
    last_metrics = None

    contents = [
        types.Content(role="user", parts=[types.Part(
            text=f"Summarize the file at '{filepath}' and save it to Notion with the title '{os.path.splitext(filename)[0]}'."
        )])
    ]

    # ── THE LOOP ──────────────────────────────
    while True:
        # "Think" -> model reads entire contents & system prompt, then responds with text and optional tool calls
        # everytime this is called, it will always read the contents so that it updated
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=contents,
            config=config, # where the system prompt is referenced, along with tool availability and response format requirements
        )

        if not response.function_calls:
            # No tool calls → model is done, write final output and move on
            output_path = os.path.join(OUTPUT_FOLDER, os.path.splitext(filename)[0] + ".md")
            with open(output_path, "w", encoding="utf-8") as out:
                out.write(response.text)
            print(f"[3] Summary saved → {output_path}")
            break

        # For retention of context -> Append the what the model said before processing any calls
        contents.append(response.candidates[0].content)

        # "ACT" -> Process ALL function calls in the turn
        tool_results = []
        for fc in response.function_calls:
            print(f"[1] Model REQUESTED: {fc.name}")

            if fc.name == "read_file":
                result = available_tools[fc.name](**fc.args)
                source_text = result
                print(f"    ↳ read {len(result)} chars")

            elif fc.name == "write_to_notion":
                summary_text = fc.args.get("summary", "")        # .get avoids KeyError if arg missing
                metrics = judge_summary_quality(source_text, summary_text)
                last_metrics = metrics                            # reuse at end — no second judge call

                faith_ok = metrics["is_faithful"]                 # HARD gate for faithfullness: zero tolerations for hallucinations
                cover_ok = metrics["coverage_score"] >= 4         # GRADED gate: mostly-complete is ok

                if faith_ok and cover_ok:
                    result = available_tools[fc.name](**fc.args)
                    print(f"    ↳ gate PASSED (faithful, coverage {metrics['coverage_score']}/5)")

                elif retry_count >= MAX_RETRIES:
                    # Deliberate blast-radius choice: do NOT write known-bad output to Notion.
                    result = (f"NOT SAVED — still failed the quality gate after {MAX_RETRIES} attempts. "
                              f"Stop and do not call any more tools.")
                    print(f"    ↳ max retries reached — NOT saved, flagged for human review")

                else:
                    retry_count += 1
                    reasons = []
                    if not faith_ok:
                        reasons.append(f"Remove these unsupported claims: {metrics['unsupported_claims']}.")
                    if not cover_ok:
                        reasons.append(f"Add these missing key points: {metrics['missing_points']}.")
                    result = (f"REJECTED: {' '.join(reasons)} "
                              f"Revise using ONLY the source, then call write_to_notion again. "
                              f"({retry_count}/{MAX_RETRIES} attempts used)")
                    print(f"    ↳ gate REJECTED — retry {retry_count}/{MAX_RETRIES}")

            else:
                result = available_tools[fc.name](**fc.args)

            # appends all tool results to var
            tool_results.append(
                types.Part.from_function_response(name=fc.name, response={"result": result})
            )

        # Observe -> append ALL tool results to the contents so the model can "observe" from them in the next turn
        print(f"[2] Tool(s) ran. Feeding back...")
        contents.append(types.Content(role="user", parts=tool_results))

    # ── POST-LOOP SCORECARD (only if we have both source and summary) ──
    if source_text and summary_text:
        print(f"\n>>> EVALUATING {filename}")
        evaluate(source_text, summary_text, metrics=last_metrics)